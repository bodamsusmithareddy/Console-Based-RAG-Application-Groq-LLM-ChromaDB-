"""
File: rag_application.py
Description: A Retrieval-Augmented Generation (RAG) Console-Based Application
Author: Susmitha Reddy Bodam
Date: December 7 2024
Version: 2.0

Purpose:
This script implements a console-based RAG application that allows users to:
1. Upload multiple .pdf or .pptx files and extract their content.
2. Index the extracted content for semantic search using ChromaDB.
3. Combine multiple indexes into a unified searchable index.
4. Query the indexed data and generate responses using the Groq LLM.

Requirements:
- Python 3.8 or higher
- Libraries: os, fitz (PyMuPDF), python-pptx, langchain_chroma, sentence_transformers, langchain_groq, concurrent.futures

Usage:
Run this script in a Python environment and follow the prompts to upload files, index their content, and query the data.
"""

import os
import fitz  # PyMuPDF
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor
from langchain_chroma import Chroma
from sentence_transformers import SentenceTransformer
from langchain_groq import ChatGroq

# Set your Groq API key
os.environ["GROQ_API_KEY"] = ""  # Replace with your actual API key

# Global constants
CHROMA_DB_DIR = "chroma_db/"
COMBINED_DB_DIR = "combined_index/"

# ------------------- File Upload and Text Extraction -------------------

def extract_text_from_pptx(file_path):
    """
    Extracts text from a PowerPoint (.pptx) file.

    Args:
        file_path (str): Path to the PowerPoint file.

    Returns:
        str: Extracted text or None if extraction fails.
    """
    try:
        presentation = Presentation(file_path)
        text_content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_content.append(paragraph.text.strip())
        return "\n".join(text_content)
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return None

def extract_text_from_pdf(file_path):
    """
    Extracts text from a PDF file.

    Args:
        file_path (str): Path to the PDF file.

    Returns:
        str: Extracted text or None if extraction fails.
    """
    try:
        document = fitz.open(file_path)
        text_content = []
        for page in document:
            blocks = page.get_text("blocks")  # Extract text as blocks
            blocks.sort(key=lambda block: (block[1], block[0]))  # Sort by vertical, then horizontal position
            for block in blocks:
                block_text = block[4].strip()
                if block_text:  # Filter out empty blocks
                    text_content.append(block_text)
        extracted_text = "\n".join(text_content)
        print("\n--- Extracted Text from PDF ---\n")
        print(extracted_text)
        print("\n--- End of Extracted Text ---\n")
        return extracted_text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return None

# ------------------- Embedding Function Wrapper -------------------

class SentenceTransformerEmbedding:
    """
    A wrapper for the SentenceTransformer model to provide embedding methods.

    Methods:
        embed_documents: Embeds a list of texts for document indexing.
        embed_query: Embeds a single query for retrieval.
    """
    def __init__(self, model_name):
        self.model = SentenceTransformer(model_name)
    
    def embed_documents(self, texts):
        embeddings = self.model.encode(texts, show_progress_bar=True)
        return [embedding.tolist() for embedding in embeddings]
    
    def embed_query(self, query):
        query_embedding = self.model.encode([query], show_progress_bar=False)
        return query_embedding[0].tolist()
# ------------------- Indexing and Retrieval -------------------

def create_index(data):
    """
    Creates a ChromaDB index for semantic search.

    Args:
        data (str): Text data to index.

    Returns:
        Chroma: A ChromaDB object for semantic search.
    """
    try:
        # Chunk data into logical sentences or paragraphs for better indexing
        chunk_size = 500
        sentences = [line.strip() for line in data.split("\n") if line.strip()]
        chunks = [" ".join(sentences[i:i + chunk_size]) for i in range(0, len(sentences), chunk_size)]
        
        embedding_function = SentenceTransformerEmbedding('all-MiniLM-L6-v2')
        if not os.path.exists(CHROMA_DB_DIR):
            os.makedirs(CHROMA_DB_DIR)
        db = Chroma(persist_directory=CHROMA_DB_DIR, embedding_function=embedding_function)
        db.add_texts(texts=chunks, metadatas=None)
        print("Documents successfully indexed.")
        return db
    except Exception as e:
        print(f"Error creating index: {e}")
        return None

def process_file(file_path):
    """
    Processes a single file to extract text and create an index.

    Args:
        file_path (str): Path to the file to be processed.

    Returns:
        Chroma: A ChromaDB object for the indexed content.
    """
    extracted_text = extract_text_from_pdf(file_path) if file_path.endswith('.pdf') else extract_text_from_pptx(file_path)
    if extracted_text:
        return create_index(extracted_text)
    return None

def combine_indexes(index_list, embedding_function):
    """
    Combines multiple indexes into a single ChromaDB index.

    Args:
        index_list (list): List of ChromaDB indexes to combine.
        embedding_function: The embedding function used by ChromaDB.

    Returns:
        Chroma: A combined ChromaDB object.
    """
    try:
        combined_db = Chroma(persist_directory=COMBINED_DB_DIR, embedding_function=embedding_function)

        for index in index_list:
            if index:
                # Retrieve all indexed data from the current index
                all_docs = index.search("", k=index._collection.count())  # Retrieve all documents
                for doc in all_docs:
                    combined_db.add_texts([doc.page_content], metadatas=[doc.metadata])

        print("Combined all indexes into a unified database.")
        return combined_db
    except Exception as e:
        print(f"Error combining indexes: {e}")
        return None

def query_index(db, query):
    """
    Queries the ChromaDB index for relevant context.

    Args:
        db (Chroma): The ChromaDB object.
        query (str): The user query.

    Returns:
        str: Retrieved context or None if no results are found.
    """
    try:
        results = db.similarity_search(query, k=3)
        if results:
            contexts = [result.page_content for result in results]
            return " ".join(contexts)
        else:
            print("No relevant context found in similarity search.")
            return None
    except Exception as e:
        print(f"Error querying index: {e}")
        return None

# ------------------- LLM Integration -------------------

def generate_response(context, question):
    """
    Generates a response using the Groq LLM.

    Args:
        context (str): Retrieved context from the database.
        question (str): User's question.

    Returns:
        str: LLM-generated response.
    """
    try:
        llm = ChatGroq(model="llama-3.1-70b-versatile")
        prompt = f"""
        Use the following retrieved context to answer the question:
        {context}
        If the context doesn't fully answer the question, indicate that explicitly.
        Question: {question}
        Answer:"""
        response = llm.invoke(prompt)
        return response.content
    except Exception as e:
        print(f"Error generating response: {e}")
        return "Sorry, I encountered an issue generating the response."

# ------------------- Main Application -------------------

def main():
    """
    Main entry point for the RAG console-based application.
    """
    print("Welcome to the Enhanced RAG Console-Based Application!")
    
    # Step 1: Batch File Processing
    file_paths = input("Enter the paths to your files (comma-separated): ").split(',')
    file_paths = [path.strip() for path in file_paths]
    
    print("Processing files...")
    with ThreadPoolExecutor() as executor:
        indexed_results = list(executor.map(process_file, file_paths))

    # Step 2: Combine Multiple Indexes
    embedding_function = SentenceTransformerEmbedding('all-MiniLM-L6-v2')
    combined_db = combine_indexes(indexed_results, embedding_function)

    if not combined_db:
        print("Failed to combine indexes. Exiting application.")
        return

    print("Combined index created successfully!")

    # Step 3: Query Loop
    print("\nYou can now start querying the combined database.")
    while True:
        query = input("\nEnter your query (or type 'exit' to quit): ")
        if query.lower() == 'exit':
            print("Goodbye!")
            break

        # Retrieve relevant context
        context = query_index(combined_db, query)
        if not context:
            print("No relevant context found. Try rephrasing your query.")
            continue

        # Generate response
        response = generate_response(context, query)
        print(f"\nResponse:\n{response}")

if __name__ == "__main__":
    main()
